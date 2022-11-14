'''
This is a diploma generator but also it can be used for any powerpoint template that needs change the name of a larger files from a excel (.csv) file.
This example have the propurse of demostrate it application:
The csv have this structure:
"Name,proyect_name(don't use in this example),proyect_number"
And in the same carpet is the powerpoint diploma template
'''
import csv
import win32com.client


def PPTtoPDF(inputFileName, outputFileName, formatType = 32):
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    powerpoint.Visible = 1

    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()

def search_and_replace(search_str, repl_str,search_str2, repl_str2, input, output):
    from pptx import Presentation
    prs = Presentation(input)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if(shape.text.find(search_str))!=-1:
                    text_frame = shape.text_frame
                    cur_text = text_frame.paragraphs[0].runs[0].text
                    new_text = cur_text.replace(str(search_str), str(repl_str))
                    text_frame.paragraphs[0].runs[0].text = new_text
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if(shape.text.find(search_str2))!=-1:
                    text_frame = shape.text_frame
                    cur_text = text_frame.paragraphs[0].runs[0].text
                    new_text = cur_text.replace(str(search_str2), str(repl_str2))
                    text_frame.paragraphs[0].runs[0].text = new_text
    prs.save(output)

with open(r'c:\Users\hurri\Documents\ProyectosPython\GeneradorDeDiplomas2_0\names_info.csv', newline='') as File:
    reader = csv.reader(File)
    line=0
    for row in reader:
        if line!=0:
            name = row[0]
            first_name = name.find(" ")
            print(row[0][0:first_name])
            print(row[2])
            if line-1>=10:
                search_and_replace('Name_tag',row[0],"tag",f"P-0{line}",r'c:\Users\hurri\Documents\ProyectosPython\GeneradorDeDiplomas2_0\participation_diploma.pptx',
                    f"Diplomas/{row[2]}-{row[0][0:first_name]}.pptx")
                PPTtoPDF(f"Diplomas/{row[2]}-{row[0][0:first_name]}.pptx",f"Diplomas/PDF/{row[2]}-{row[0][0:first_name]}.pdf")

            else:
                search_and_replace('Name_tag', row[0], "tag", f"P-00{line}",
                                   r'c:\Users\hurri\Documents\ProyectosPython\GeneradorDeDiplomas2_0\participation_diploma.pptx',
                                   f"c:/Users/hurri/Documents/ProyectosPython/GeneradorDeDiplomas2_0/Diplomas/{row[2]}-{row[0][0:first_name]}.pptx")
                PPTtoPDF(f"c:/Users/hurri/Documents/ProyectosPython/GeneradorDeDiplomas2_0/Diplomas/{row[2]}-{row[0][0:first_name]}.pptx",f"c:\\Users\\hurri\\Documents\\ProyectosPython\\GeneradorDeDiplomas2_0\\Diplomas\\PDF\\{row[2]}-{row[0][0:first_name]}.pdf")
        line=line+1